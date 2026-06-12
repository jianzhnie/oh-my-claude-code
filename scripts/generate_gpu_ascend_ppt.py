#!/usr/bin/env python3
"""Generate a 10-slide PowerPoint from the GPU vs Ascend ecosystem research report."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Color palette
NAVY = RGBColor(0x1A, 0x23, 0x7E)
NVIDIA_GREEN = RGBColor(0x76, 0xB9, 0x00)
ASCEND_RED = RGBColor(0xCF, 0x0A, 0x2C)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x00, 0x6C, 0xFF)

OUTPUT_PATH = (
    "/Users/robin/work_dir/oh-my-claude-code/docs/gpu-ascend-ecosystem-report.pptx"
)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a colored title bar at the top."""
    # Title bar background
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.35)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p2.font.name = "Microsoft YaHei"
        p2.space_before = Pt(4)

    return Inches(1.55)


def add_footer(slide, page_num, total):
    """Add footer with page number."""
    footer = slide.shapes.add_textbox(
        Inches(11.5), Inches(7.35), Inches(1.5), Inches(0.3)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.RIGHT


def add_bullet_box(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    title=None,
    title_color=NAVY,
    bullet_color=DARK_GRAY,
    font_size=16,
):
    """Add a bullet point text box."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(10)

    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if title or i > 0 else tf.paragraphs[0]
        p.text = f"• {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = bullet_color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(8)
        p.level = 0

    return box


def add_two_column_compare(
    slide, top, gpu_title, gpu_items, ascend_title, ascend_items
):
    """Add a two-column comparison layout."""
    # NVIDIA column
    col1_left = Inches(0.5)
    col_width = Inches(6.0)
    col_height = Inches(5.0)

    # NVIDIA header
    header1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, col1_left, top, col_width, Inches(0.5)
    )
    header1.fill.solid()
    header1.fill.fore_color.rgb = NVIDIA_GREEN
    header1.line.fill.background()
    h1_tf = header1.text_frame
    h1_tf.paragraphs[0].text = gpu_title
    h1_tf.paragraphs[0].font.size = Pt(20)
    h1_tf.paragraphs[0].font.bold = True
    h1_tf.paragraphs[0].font.color.rgb = WHITE
    h1_tf.paragraphs[0].font.name = "Microsoft YaHei"
    h1_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # NVIDIA content
    add_bullet_box(
        slide,
        col1_left,
        top + Inches(0.6),
        col_width,
        col_height,
        gpu_items,
        title=None,
        bullet_color=DARK_GRAY,
        font_size=15,
    )

    # Ascend column
    col2_left = Inches(6.833)
    header2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, col2_left, top, col_width, Inches(0.5)
    )
    header2.fill.solid()
    header2.fill.fore_color.rgb = ASCEND_RED
    header2.line.fill.background()
    h2_tf = header2.text_frame
    h2_tf.paragraphs[0].text = ascend_title
    h2_tf.paragraphs[0].font.size = Pt(20)
    h2_tf.paragraphs[0].font.bold = True
    h2_tf.paragraphs[0].font.color.rgb = WHITE
    h2_tf.paragraphs[0].font.name = "Microsoft YaHei"
    h2_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Ascend content
    add_bullet_box(
        slide,
        col2_left,
        top + Inches(0.6),
        col_width,
        col_height,
        ascend_items,
        title=None,
        bullet_color=DARK_GRAY,
        font_size=15,
    )


def add_table(
    slide,
    left,
    top,
    width,
    height,
    headers,
    rows,
    header_color=NAVY,
    header_text_color=WHITE,
):
    """Add a styled table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Set column widths
    col_width = width / num_cols
    for col in table.columns:
        col.width = int(col_width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = header_text_color
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: Title =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Decorative shapes
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10.5), Inches(-1), Inches(4), Inches(4)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0x30, 0x3F, 0x9F)
    shape1.line.fill.background()

    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-1.5), Inches(5), Inches(4), Inches(4)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0x13, 0x1A, 0x5E)
    shape2.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(11.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GPU 与 Ascend 大模型训练推理生态调研"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.0), Inches(11.333), Inches(0.8)
    )
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "NVIDIA GPU vs 华为昇腾 NPU 生态全景、训练推理实践与选型建议"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER

    # Date and author
    info_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.8), Inches(11.333), Inches(0.6)
    )
    tf3 = info_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "调研日期：2026/06/12"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p3.font.name = "Microsoft YaHei"
    p3.alignment = PP_ALIGN.CENTER

    add_footer(slide, 1, 10)

    # ===== Slide 2: 生态全景图 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "一、生态全景图", "从芯片到应用的全栈分层对比")
    add_footer(slide, 2, 10)

    headers = ["层级", "NVIDIA GPU 生态", "华为 Ascend 生态"]
    rows = [
        ["训练框架", "PyTorch/JAX + CUDA", "MindSpore + torch_npu"],
        ["分布式框架", "Megatron-LM / DeepSpeed / FSDP", "MindSpeed-LLM / MM / RL"],
        ["推理引擎", "vLLM / TensorRT-LLM / SGLang", "vLLM-Ascend / MindIE / SGLang"],
        ["通信库", "NCCL", "HCCL"],
        ["算子开发", "CUDA / Triton / CUTLASS", "Ascend C / TBE / Ascend-Triton"],
        ["硬件平台", "H100 / H200 / B200 / GB200", "910B / 910C / 950PR / 950DT"],
        ["互联方案", "NVLink + NVSwitch + IB/RoCE", "HCCS + MatrixLink / CloudMatrix"],
    ]
    add_table(
        slide,
        Inches(0.5),
        Inches(1.8),
        Inches(12.333),
        Inches(5.0),
        headers,
        rows,
        header_color=NAVY,
    )

    # ===== Slide 3: 训练硬件对比 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(
        slide, "二、训练硬件对比", "单卡算力与显存：NVIDIA 领先，昇腾超节点补差距"
    )
    add_footer(slide, 3, 10)

    headers = ["指标", "H100 SXM5", "B200 SXM6", "昇腾 910B", "昇腾 910C", "昇腾 950DT"]
    rows = [
        [
            "BF16 算力",
            "989 TFLOPS",
            "~4,500 TFLOPS",
            "~376 TFLOPS",
            "~780 TFLOPS",
            "~1,000 TFLOPS",
        ],
        ["显存容量", "80 GB", "192 GB", "64 GB", "升级", "144 GB"],
        ["显存带宽", "3.35 TB/s", "8.0 TB/s", "~1.6 TB/s", "3.2 TB/s", "4 TB/s"],
        ["互联带宽", "900 GB/s", "1.8 TB/s", "~392 GB/s", "784 GB/s", "2 TB/s"],
        ["单卡功耗", "700 W", "1,000 W", "310–400 W", "~400 W", "~550 W"],
        ["定位", "训练旗舰", "训练/推理", "训练", "训练/推理", "训练/解码"],
    ]
    add_table(
        slide,
        Inches(0.3),
        Inches(1.75),
        Inches(12.733),
        Inches(4.8),
        headers,
        rows,
        header_color=NAVY,
    )

    # Key insight
    insight = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(6.6),
        Inches(12.333),
        Inches(0.55),
    )
    insight.fill.solid()
    insight.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    insight.line.fill.background()
    itf = insight.text_frame
    itf.paragraphs[
        0
    ].text = "关键结论：单卡 NVIDIA 仍领先，但 CloudMatrix 384 超节点通过系统级设计实现算力/内存/互联的集群级超越。"
    itf.paragraphs[0].font.size = Pt(14)
    itf.paragraphs[0].font.color.rgb = NAVY
    itf.paragraphs[0].font.name = "Microsoft YaHei"
    itf.paragraphs[0].alignment = PP_ALIGN.CENTER
    insight.text_frame.paragraphs[0].font.bold = True

    # ===== Slide 4: 训练软件栈与并行策略 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(
        slide, "三、训练软件栈与并行策略", "Ascend 已覆盖主流并行策略，生态快速补齐"
    )
    add_footer(slide, 4, 10)

    add_two_column_compare(
        slide,
        Inches(1.75),
        "GPU 生态",
        [
            "框架：PyTorch / JAX / TensorFlow 原生 CUDA",
            "分布式：Megatron-LM / DeepSpeed / FSDP",
            "通信：NCCL 2.21+",
            "算子：cuBLAS / cuDNN / FlashAttention-3",
            "成熟度高，第三方生态丰富",
        ],
        "Ascend 生态",
        [
            "框架：MindSpore + PyTorch via torch_npu",
            "分布式：MindSpeed-LLM / MindSpeed-MM / MindSpeed-RL",
            "通信：HCCL",
            "算子：AscendCL / Ascend FlashAttention / AMLA",
            "2025–2026 年框架/算子快速追赶",
        ],
    )

    # Parallel strategies box
    para_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(5.6),
        Inches(12.333),
        Inches(1.1),
    )
    para_box.fill.solid()
    para_box.fill.fore_color.rgb = LIGHT_GRAY
    para_box.line.fill.background()
    ptf = para_box.text_frame
    ptf.paragraphs[
        0
    ].text = "并行策略全覆盖：DP / TP / PP / SP / EP / CP，昇腾 MindSpeed-LLM 支持 3D/4D 并行组合，MoE 场景支持 EPLB 负载均衡。"
    ptf.paragraphs[0].font.size = Pt(16)
    ptf.paragraphs[0].font.color.rgb = DARK_GRAY
    ptf.paragraphs[0].font.name = "Microsoft YaHei"
    ptf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ptf.paragraphs[0].font.bold = True

    # ===== Slide 5: 超大规模模型训练实践 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(
        slide,
        "四、超大规模模型在 Ascend NPU 上的训练实践",
        "从千卡到万卡，国产算力已支撑万亿参数模型",
    )
    add_footer(slide, 5, 10)

    headers = ["模型", "参数规模", "训练集群", "关键指标", "意义"]
    rows = [
        [
            "盘古 Ultra",
            "135B Dense",
            "8,192 张昇腾 NPU",
            "MFU 50%",
            "首个纯昇腾稠密大模型",
        ],
        [
            "盘古 Ultra MoE",
            "718B (39B)",
            "6,000+ / 万卡",
            "MFU 41%",
            "准万亿 MoE，媲美 DeepSeek-R1",
        ],
        [
            "DeepSeek-V4",
            "Pro 1.6T (49B)",
            "Ascend 910B 千卡",
            "CANN 重写 200+ 算子",
            "首个公开全栈适配万亿级模型",
        ],
        [
            "LongCat-2.0",
            "万亿+",
            "5–6 万张昇腾",
            "1M 上下文",
            "唯一公开确认国产算力万亿参数预训练",
        ],
        [
            "讯飞星火 X2-Flash",
            "30B MoE",
            "昇腾 910B 集群",
            "256K 上下文",
            "国产算力 Agent 时代模型",
        ],
        [
            "GLM-Image",
            "多模态",
            "Atlas 800T A2 + MindSpore",
            "SOTA 多模态",
            "首个国产算力全程训练",
        ],
    ]
    add_table(
        slide,
        Inches(0.25),
        Inches(1.75),
        Inches(12.833),
        Inches(4.9),
        headers,
        rows,
        header_color=ASCEND_RED,
    )

    # ===== Slide 6: 推理生态对比 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(
        slide, "五、大模型推理生态对比", "昇腾与 NVIDIA 差距最小的场景在推理侧"
    )
    add_footer(slide, 6, 10)

    add_two_column_compare(
        slide,
        Inches(1.75),
        "GPU 推理",
        [
            "引擎：vLLM / TensorRT-LLM / SGLang",
            "量化：FP8 / FP4 / INT8 / AWQ / GPTQ",
            "Attention：FlashAttention-2/3",
            "部署：Triton / NVIDIA NIM",
            "单卡显存：H200 141GB / B200 192GB",
        ],
        "Ascend 推理",
        [
            "引擎：vLLM-Ascend / MindIE / SGLang",
            "量化：INT8 / FP8 / C8 INT8 KV Cache",
            "Attention：Ascend FlashAttention / AMLA",
            "部署：MindIE-Service / ModelArts",
            "超节点：CloudMatrix 384 大显存集群",
        ],
    )

    # DeepSeek-R1 inference performance
    perf_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(5.5),
        Inches(12.333),
        Inches(1.2),
    )
    perf_box.fill.solid()
    perf_box.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    perf_box.line.fill.background()
    ptf = perf_box.text_frame
    ptf.paragraphs[0].text = "DeepSeek-R1 671B 推理实测（50ms 时延约束）"
    ptf.paragraphs[0].font.size = Pt(16)
    ptf.paragraphs[0].font.bold = True
    ptf.paragraphs[0].font.color.rgb = ASCEND_RED
    ptf.paragraphs[0].font.name = "Microsoft YaHei"
    ptf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p2 = ptf.add_paragraph()
    p2.text = (
        "H100：~45 ms TTFT  |  昇腾 910B：~52 ms TTFT，单卡 decode 吞吐 1,920 tokens/s"
    )
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_GRAY
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER

    # ===== Slide 7: 关键融合算子 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(
        slide,
        "六、关键融合算子与加速库",
        "Ascend 在 Attention / MoE / MLA 上从可用走向好用",
    )
    add_footer(slide, 7, 10)

    # Three boxes
    box_width = Inches(4.0)
    box_height = Inches(4.6)
    gap = Inches(0.166)
    start_left = Inches(0.5)
    tops = [Inches(1.75)] * 3
    lefts = [
        start_left,
        start_left + box_width + gap,
        start_left + 2 * (box_width + gap),
    ]

    boxes = [
        (
            "FlashAttention / MLA",
            [
                "GPU：FlashAttention-2/3 / FlashMLA (66.7%)",
                "Ascend：Ascend FlashAttention / AMLA (86.8%)",
                "AMLA 面向 DeepSeek MLA 优化，利用率超越 FlashMLA",
            ],
            ACCENT_BLUE,
        ),
        (
            "MoE 通信",
            [
                "GPU：DeepEP (sub-150μs)",
                "Ascend：DeepEP-Ascend / ascend_fuseep",
                "兼容 DeepEP API，低延迟模式已支持 DeepSeek-V3",
            ],
            NVIDIA_GREEN,
        ),
        (
            "算子开发",
            [
                "GPU：CUDA / Triton / CUTLASS",
                "Ascend：Ascend C / TBE / TileLang-Ascend",
                "TileLang-Ascend 2026 年发布 FlashAttention、DeepSeek-V4 kernel",
            ],
            ASCEND_RED,
        ),
    ]

    for i, (title, bullets, color) in enumerate(boxes):
        # Box background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, lefts[i], tops[i], box_width, box_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.color.rgb = color
        bg.line.width = Pt(2)

        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, lefts[i], tops[i], box_width, Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        htf = header.text_frame
        htf.paragraphs[0].text = title
        htf.paragraphs[0].font.size = Pt(16)
        htf.paragraphs[0].font.bold = True
        htf.paragraphs[0].font.color.rgb = WHITE
        htf.paragraphs[0].font.name = "Microsoft YaHei"
        htf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Bullets
        add_bullet_box(
            slide,
            lefts[i] + Inches(0.15),
            tops[i] + Inches(0.65),
            box_width - Inches(0.3),
            box_height - Inches(0.8),
            bullets,
            title=None,
            bullet_color=DARK_GRAY,
            font_size=13,
        )

    # ===== Slide 8: 社区框架支持 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(
        slide,
        "七、社区框架与工具对 Ascend NPU 的支持",
        "2024–2026 年主流社区框架密集适配昇腾",
    )
    add_footer(slide, 8, 10)

    headers = ["场景", "推荐框架", "状态", "备注"]
    rows = [
        ["快速 SFT/LoRA", "LLaMA-Factory / ms-swift", "✅ 官方支持", "易用、文档完善"],
        [
            "大规模预训练",
            "MindSpeed-LLM / Megatron + torch_npu",
            "✅ 官方优化",
            "并行策略最全",
        ],
        [
            "RLHF / GRPO",
            "MindSpeed RL / verl / ms-swift",
            "✅ 可用",
            "MindSpeed RL 性能最强",
        ],
        ["高吞吐推理", "vLLM-Ascend", "✅ 事实标准", "模型覆盖最广"],
        ["结构化生成", "SGLang", "🚧 快速补齐", "2026 Q2 重点"],
        ["政企合规推理", "MindIE", "✅ 生产级", "华为自研原生"],
        [
            "HuggingFace 用户",
            "transformers + accelerate + trl",
            "✅ 可用",
            "零迁移成本",
        ],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.75),
        Inches(12.533),
        Inches(4.9),
        headers,
        rows,
        header_color=NAVY,
    )

    # ===== Slide 9: 工程实践与 TCO =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "八、工程实践与 TCO", "迁移、版本、成本、部署的关键落地信息")
    add_footer(slide, 9, 10)

    add_two_column_compare(
        slide,
        Inches(1.75),
        "迁移与版本",
        [
            "自动迁移：transfer_to_npu 一行注入",
            "手动迁移：torch.cuda→torch.npu，nccl→hccl",
            "自定义 CUDA kernel 需重写为 Ascend C",
            "典型周期：简单模型 1–2 周，生产服务 1–3 月",
            "版本严格耦合：CANN/PyTorch/torch_npu/vLLM-Ascend",
        ],
        "成本与部署",
        [
            "8×H100 云租赁：$8–98/小时（spot vs Azure）",
            "8×910B 云租赁：¥20–45/小时",
            "昇腾采购价显著低于 H100，国产化补贴后更优",
            "K8s：vLLM-Ascend / MindIE Service 官方方案",
            "PD 分离、KV Cache 量化已支持生产部署",
        ],
    )

    # ===== Slide 10: 结论与趋势 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_footer(slide, 10, 10)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.9)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "结论与趋势"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

    # Conclusion bullets
    bullets = [
        "极致性能、最新算法、全球化部署：NVIDIA GPU 仍是唯一选择",
        "国产化、政策合规、供应链安全、成本可控：Ascend 已进入「好用」阶段，推理侧值得投入",
        "软件栈收敛：PyTorch(torch_npu) + vLLM-Ascend + MindSpeed-LLM + CANN/HCCL 成为主路径",
        "关键算子快速追赶：DeepEP-Ascend、AMLA、TileLang-Ascend 是 2025–2026 关键里程碑",
        "集群架构成为新焦点：从单卡竞赛转向超节点级系统设计",
    ]
    add_bullet_box(
        slide,
        Inches(1.0),
        Inches(1.6),
        Inches(11.333),
        Inches(4.5),
        bullets,
        title=None,
        bullet_color=WHITE,
        font_size=22,
    )

    # Final statement
    final = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.8)
    )
    ftf = final.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "掌握 PyTorch + vLLM 即可同时覆盖 GPU 与 Ascend；Ascend 特有技能集中在 CANN / MindSpeed / HCCL 调优。"
    fp.font.size = Pt(20)
    fp.font.bold = True
    fp.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)
    fp.font.name = "Microsoft YaHei"
    fp.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT_PATH)
    print(f"PPT saved to: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_presentation()
