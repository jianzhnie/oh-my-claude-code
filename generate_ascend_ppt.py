#!/usr/bin/env python3
"""Generate a polished Ascend ecosystem and practice PPT."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x0F, 0x1F, 0x4B)  # deep navy
ACCENT = RGBColor(0xF5, 0xA6, 0x23)  # warm gold
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xE2)  # soft blue
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)  # near black
MUTED_TEXT = RGBColor(0x66, 0x66, 0x66)  # grey
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)  # subtle grey-blue background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _add_shape(slide, shape_type, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    font_color=DARK_TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Microsoft YaHei",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tb


def _add_footer(slide, page_num, total):
    # Bottom accent bar
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        SLIDE_HEIGHT - Inches(0.18),
        SLIDE_WIDTH,
        Inches(0.18),
        fill=PRIMARY,
    )
    # Page number
    _add_textbox(
        slide,
        SLIDE_WIDTH - Inches(1.2),
        SLIDE_HEIGHT - Inches(0.55),
        Inches(1),
        Inches(0.3),
        f"{page_num} / {total}",
        font_size=12,
        font_color=MUTED_TEXT,
        align=PP_ALIGN.RIGHT,
    )


def _add_section_badge(slide, text, top=None):
    if top is None:
        top = Inches(0.35)
    # Small gold badge at top-left
    _add_textbox(
        slide,
        Inches(0.5),
        top,
        Inches(4),
        Inches(0.35),
        text,
        font_size=14,
        font_color=ACCENT,
        bold=True,
    )
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        top + Inches(0.32),
        Inches(0.8),
        Inches(0.04),
        fill=ACCENT,
    )


def _set_slide_title(slide, title_text):
    # Title placeholder already exists on title-content layout
    title = slide.shapes.title
    title.text = title_text
    tf = title.text_frame
    tf.paragraphs[0].font.size = Pt(34)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    return title


def _add_bullets(
    slide,
    bullets,
    left,
    top,
    width,
    height,
    font_size=19,
    numbered=False,
    accent=ACCENT,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        marker = f"{i + 1}." if numbered else "•"
        p.text = f"{marker} {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(14)
        p.level = 0
    return tb


def _add_card(
    slide, left, top, width, height, title, bullets, title_color=PRIMARY, bullet_size=17
):
    # Card background
    _add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill=WHITE,
        line=RGBColor(0xE0, 0xE4, 0xEA),
    )
    # Title
    _add_textbox(
        slide,
        left + Inches(0.25),
        top + Inches(0.2),
        width - Inches(0.5),
        Inches(0.5),
        title,
        font_size=22,
        font_color=title_color,
        bold=True,
    )
    # Accent underline
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.25),
        top + Inches(0.62),
        Inches(0.6),
        Inches(0.05),
        fill=ACCENT,
    )
    # Bullets
    _add_bullets(
        slide,
        bullets,
        left + Inches(0.25),
        top + Inches(0.8),
        width - Inches(0.5),
        height - Inches(1.1),
        font_size=bullet_size,
    )


def _add_table(
    slide,
    data,
    left,
    top,
    width,
    height,
    header_fill=PRIMARY,
    header_text=WHITE,
    body_text=DARK_TEXT,
):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for i, row in enumerate(data):
        for j, text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = text
            cell.text_frame.word_wrap = True
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(15)
            p.font.name = "Microsoft YaHei"
            p.font.color.rgb = header_text if i == 0 else body_text
            p.font.bold = i == 0
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table


# ---------------------------------------------------------------------------
# Content slides
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # We know final page count is 14; used for footer.
    TOTAL_PAGES = 14

    # ------------------------------------------------------------------
    # Slide 1: Cover
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background gradient-like via full-size navy rect
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        SLIDE_HEIGHT,
        fill=PRIMARY,
    )
    # Decorative gold stripe
    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.8),
        Inches(3.2),
        Inches(1.2),
        Inches(0.08),
        fill=ACCENT,
    )
    # Title
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(11.5),
        Inches(1.2),
        "昇腾生态与大模型训练实践",
        font_size=54,
        font_color=WHITE,
        bold=True,
    )
    # Subtitle
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(3.5),
        Inches(11.5),
        Inches(0.7),
        "GPU vs Ascend 生态调研  ·  鹏城实验室昇腾训练实践",
        font_size=26,
        font_color=ACCENT,
    )
    # Date / org
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(5.8),
        Inches(11.5),
        Inches(0.5),
        "2026 年 6 月",
        font_size=18,
        font_color=RGBColor(0xAA, 0xB2, 0xC5),
    )
    _add_footer(slide, 1, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 2: Ecosystem Overview
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "一、昇腾生态全景：分层对比")
    _add_section_badge(slide, "ECOSYSTEM OVERVIEW")
    data = [
        ["层级", "GPU 生态 (NVIDIA)", "Ascend 生态 (华为)", "说明"],
        [
            "应用层",
            "GPT、Claude、Llama",
            "盘古、DeepSeek-V4 适配、LongCat",
            "模型与硬件解耦",
        ],
        [
            "训练框架",
            "PyTorch、JAX、Megatron",
            "PyTorch + MindSpeed-LLM",
            "torch_npu 兼容",
        ],
        [
            "推理引擎",
            "vLLM、TensorRT-LLM",
            "vLLM-Ascend、MindIE",
            "vLLM-Ascend 为主入口",
        ],
        ["算子/通信", "CUDA、NCCL", "Ascend C、HCCL", "自定义 kernel 与集合通信"],
        [
            "运行时",
            "CUDA Runtime/Graph",
            "CANN (GE + MindIR + AOE)",
            "图编译与自动调优",
        ],
        ["硬件", "H100/H200/B200/GB200", "910B/910C/950PR/950DT", "超节点弥补单卡差距"],
    ]
    _add_table(slide, data, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4))
    _add_footer(slide, 2, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 3: Hardware Comparison
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "二、训练硬件对比：NVIDIA vs Ascend")
    _add_section_badge(slide, "TRAINING HARDWARE")
    data = [
        ["维度", "NVIDIA GPU", "Ascend NPU"],
        ["主力训练芯片", "H100 / H200 / B200", "910B / 910C / 950DT"],
        [
            "峰值 BF16",
            "H100: 989 TFLOPS; B200: ~4500 TFLOPS",
            "910B: ~376 TFLOPS; 910C: ~780 TFLOPS",
        ],
        [
            "显存容量",
            "H100 80GB → H200 141GB → B200 192GB",
            "910B 64GB → 910C 128GB → 950DT 144GB",
        ],
        [
            "显存带宽",
            "H100 3.35 TB/s → B200 8.0 TB/s",
            "910B ~1.6 TB/s → 910C 3.2 TB/s → 950DT 4 TB/s",
        ],
        [
            "互联带宽",
            "NVLink 4: 900 GB/s; NVLink 5: 1.8 TB/s",
            "910B HCCS 392 GB/s 双向; 910C UB 392 GB/s 单向",
        ],
        ["关键差异", "单卡算力、显存、互联成熟", "超节点架构 + 国产化供应弥补单卡差距"],
    ]
    _add_table(slide, data, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4))
    _add_footer(slide, 3, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 4: Software Stack (two-column cards)
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "三、训练软件栈对比")
    _add_section_badge(slide, "SOFTWARE STACK")
    _add_card(
        slide,
        Inches(0.6),
        Inches(1.3),
        Inches(5.9),
        Inches(5.6),
        "GPU 生态",
        [
            "框架：PyTorch、JAX、TensorFlow",
            "分布式：Megatron-LM、DeepSpeed、FSDP",
            "通信库：NCCL 2.21+",
            "算子库：cuBLAS、cuDNN、CUTLASS",
            "Kernel 开发：CUDA、Triton",
            "Profiler：Nsight、PyTorch Profiler",
        ],
        title_color=LIGHT_BLUE,
    )
    _add_card(
        slide,
        Inches(6.8),
        Inches(1.3),
        Inches(5.9),
        Inches(5.6),
        "Ascend 生态",
        [
            "框架：PyTorch via torch_npu + MindSpore",
            "分布式：MindSpeed-LLM / -MM / -RL",
            "通信库：HCCL",
            "算子库：AscendCL、ATB",
            "Kernel 开发：Ascend C、TBE、Ascend-Triton",
            "Profiler：MindStudio Insight / Profiler",
        ],
        title_color=PRIMARY,
    )
    _add_footer(slide, 4, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 5: Parallel Strategies (numbered cards)
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "四、分布式并行策略")
    _add_section_badge(slide, "PARALLELISM STRATEGY")
    strategies = [
        ("TP 张量并行", "将 Transformer 层参数切分到多张卡，降低单卡显存占用。"),
        ("PP 流水线并行", "将模型按层切分到不同节点，扩展支持更深网络。"),
        ("EP 专家并行", "MoE 场景下将不同专家分配到不同卡/节点。"),
        ("CP 序列并行", "对长序列输入进行切分，降低激活值显存占用。"),
    ]
    positions = [(0.5, 1.35), (6.8, 1.35), (0.5, 3.95), (6.8, 3.95)]
    for (x, y), (name, desc) in zip(positions, strategies, strict=False):
        _add_card(
            slide,
            Inches(x),
            Inches(y),
            Inches(6.0),
            Inches(2.3),
            name,
            [desc],
            title_color=PRIMARY,
            bullet_size=18,
        )
    # Bottom highlight
    _add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(6.45),
        Inches(12.1),
        Inches(0.7),
        fill=RGBColor(0xE8, 0xF0, 0xFE),
    )
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(6.55),
        Inches(11.5),
        Inches(0.5),
        "典型组合：PCL-LLM-Model 采用 TP=2 / PP=8 / EP=64 / CP=1，在 2048 卡集群完成万亿参数训练。",
        font_size=17,
        font_color=PRIMARY,
        bold=True,
    )
    _add_footer(slide, 5, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 6: Large-scale Training Practice
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "五、超大规模模型在昇腾上的训练实践")
    _add_section_badge(slide, "LARGE-SCALE TRAINING")
    data = [
        ["模型", "参数规模", "训练集群", "关键指标"],
        ["盘古 Ultra", "135B Dense", "8192 张昇腾 NPU", "MFU 52%；13.2T tokens"],
        ["盘古 Ultra MoE", "718B（激活 39B）", "万卡级昇腾集群", "MFU 41%（万卡）"],
        [
            "盘古 Pro MoE",
            "72B（激活 16B）",
            "4000 颗昇腾 NPU",
            "13T tokens；三阶段预训练",
        ],
        [
            "DeepSeek-V4",
            "1.6T（激活 49B）",
            "昇腾平台适配",
            "2026.04 发布；昇腾适配验证",
        ],
        ["LongCat-2.0", "1T+", "5–6 万张国产加速卡", "1M 上下文；芯片细节未披露"],
        [
            "讯飞星火 X2-Flash",
            "30B MoE",
            "昇腾 910B 集群",
            "训练效率达同规模 A800 的 90%",
        ],
    ]
    _add_table(slide, data, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
    _add_footer(slide, 6, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 7: Training Optimization
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "六、训练优化技术与关键挑战")
    _add_section_badge(slide, "OPTIMIZATION & CHALLENGES")
    bullets = [
        "算子优化：针对昇腾 NPU Cube/Vector 单元优化 Attention、MoE Router、FFN；CANN 图编译与算子融合。",
        "通信优化：节点内 HCCS/UB 高带宽互联；节点间高效 All-to-All / All-Reduce。",
        "显存优化：激活重计算、梯度检查点、ZeRO 分片、CP 长序列并行。",
        "长稳训练：异步 Checkpoint、故障自动恢复、异常节点隔离。",
        "关键挑战：MoE All-to-All 通信、FP8 极致优化、超长稳训练稳定性、大规模 RL 训练。",
    ]
    # Left icon-like bullets
    _add_bullets(
        slide,
        bullets,
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.4),
        numbered=True,
        font_size=20,
    )
    _add_footer(slide, 7, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 8: Inference Ecosystem
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "七、推理生态对比")
    _add_section_badge(slide, "INFERENCE ECOSYSTEM")
    _add_card(
        slide,
        Inches(0.6),
        Inches(1.3),
        Inches(5.9),
        Inches(5.6),
        "GPU 推理",
        [
            "引擎：vLLM、TensorRT-LLM、SGLang",
            "量化：FP16/INT8/FP8/FP4/AWQ/GPTQ",
            "服务化：Triton Inference Server",
            "PD 分离：vLLM / SGLang",
            "优势：显存大、量化成熟、吞吐极限高",
        ],
        title_color=LIGHT_BLUE,
    )
    _add_card(
        slide,
        Inches(6.8),
        Inches(1.3),
        Inches(5.9),
        Inches(5.6),
        "Ascend 推理",
        [
            "引擎：vLLM-Ascend、MindIE、SGLang（适配中）",
            "量化：INT8/FP8 为主，C8 INT8 KV Cache 已支持",
            "服务化：MindIE-Service、ModelArts",
            "PD 分离：vLLM-Ascend / MindIE",
            "优势：与 vLLM 接口一致，迁移成本低",
        ],
        title_color=PRIMARY,
    )
    _add_footer(slide, 8, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 9: Key Fused Operators
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "八、关键融合算子与加速库")
    _add_section_badge(slide, "FUSED OPERATORS")
    bullets = [
        "FlashAttention / MLA：Ascend FlashAttention 已可用；AMLA 在昇腾 910 上达到 86.8% FLOPS 利用率。",
        "MoE 通信优化：DeepEP-Ascend 实现 sub-150μs 低延迟；EPLB 专家负载均衡。",
        "通用融合算子：Fused RMSNorm、RoPE、SwiGLU、Matmul+AllReduce+Add+RMSNorm。",
        "算子开发：Ascend C、TBE、Ascend-Triton、TileLang-Ascend。",
        "结论：Ascend 在 FlashAttention、MoE EP、MLA 等关键算子上已从'可用'走向'好用'。",
    ]
    _add_bullets(
        slide,
        bullets,
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.4),
        numbered=True,
        font_size=20,
    )
    _add_footer(slide, 9, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 10: Community Frameworks
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "九、社区框架对 Ascend 的支持")
    _add_section_badge(slide, "COMMUNITY FRAMEWORKS")
    data = [
        ["需求场景", "推荐框架（Ascend）", "备注"],
        ["快速 SFT/LoRA", "LLaMA-Factory、ms-swift", "易用、文档完善"],
        ["大规模预训练", "MindSpeed-LLM、Megatron + torch_npu", "华为官方优化最全"],
        ["RLHF / GRPO / PPO", "MindSpeed RL、verl、ms-swift", "MindSpeed RL 性能领先"],
        ["在线推理", "vLLM-Ascend", "事实标准，模型覆盖广"],
        ["结构化生成", "SGLang（NPU 适配中）", "2026 Q2 功能快速补齐"],
        ["政企合规", "MindIE", "华为自研，原生支持"],
    ]
    _add_table(slide, data, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
    _add_footer(slide, 10, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 11: Trends and Conclusion
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "十、趋势判断与结论")
    _add_section_badge(slide, "TRENDS & CONCLUSION")
    bullets = [
        "训练端：NVIDIA 仍是全球主流，昇腾通过 MindSpeed-LLM + 超节点 + 国产化政策在政企市场快速渗透。",
        "推理端：昇腾与 NVIDIA 差距最小，vLLM-Ascend 让开发者可用同一套接口服务。",
        "社区生态：LLaMA-Factory、ms-swift、verl、SGLang 等主流框架 2024–2026 年密集适配 Ascend。",
        "关键算子：DeepEP-Ascend、AMLA、TileLang-Ascend 标志着 Ascend 核心算子从可用走向好用。",
        "结论：追求极致性能选 NVIDIA；看重国产化/合规/供应链安全选 Ascend，2025–2026 年已进入'好用'阶段。",
    ]
    _add_bullets(
        slide,
        bullets,
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.4),
        numbered=True,
        font_size=20,
    )
    _add_footer(slide, 11, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 12: Practice Overview
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "十一、昇腾大模型训练实践：概览")
    _add_section_badge(slide, "PRACTICE OVERVIEW")
    bullets = [
        "7B GRPO 算法验证：基于 OpenRLHF 与 MindSpeed-RL，在 GSM8K、Math500、AIME24 上验证算法正确性。",
        "32B 模型 SFT：PCL-Reasoner-V1 在 AIME24 达到 85.7%，AIME25 达到 84.2%。",
        "32B 模型 Offline RL：PCL-Reasoner-V1.5 在 AIME24 达到 90.8%，AIME25 达到 85.7%，32B 级别榜单第一。",
        "千卡/万卡预训练：Llama3-405B MFU 43.62%；自研万亿参数 MoE 模型 PCL-LLM-Model 完成 5T tokens、30 天预训练。",
        "评估基础设施：开源 LLMEval 框架，支撑训练、蒸馏、评测全流程。",
    ]
    _add_bullets(
        slide,
        bullets,
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.4),
        numbered=True,
        font_size=20,
    )
    _add_footer(slide, 12, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 13: PCL-Reasoner & Trillion-parameter MoE
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "十二、PCL-Reasoner 与万亿参数 MoE")
    _add_section_badge(slide, "PCL-REASONER & MoE")
    _add_card(
        slide,
        Inches(0.6),
        Inches(1.3),
        Inches(5.9),
        Inches(5.6),
        "PCL-Reasoner 系列（32B）",
        [
            "V1：Qwen2.5-32B-Base + R1-0528 蒸馏数据 SFT",
            "V1.5：V1 + 自研 Offline RL",
            "AIME24：V1 85.7% → V1.5 90.8%",
            "AIME25：V1 84.2% → V1.5 85.7%",
            "已开源：HuggingFace、ModelScope、启智社区、GitHub",
        ],
        title_color=LIGHT_BLUE,
    )
    _add_card(
        slide,
        Inches(6.8),
        Inches(1.3),
        Inches(5.9),
        Inches(5.6),
        "PCL-LLM-Model（万亿 MoE）",
        [
            "架构：MoE + GQA；128 路由专家 + 1 共享专家",
            "规模：隐藏维度 7168，32 层，词表 163840",
            "集群：2048 卡（4096 Die）",
            "数据：5T tokens；训练 30 天",
            "并行：TP=2 / PP=8 / EP=64 / CP=1",
        ],
        title_color=PRIMARY,
    )
    _add_footer(slide, 13, TOTAL_PAGES)

    # ------------------------------------------------------------------
    # Slide 14: LongCat & Open Source
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title(slide, "十三、LongCat 扩展与开源贡献")
    _add_section_badge(slide, "LONGCAT & OPEN SOURCE")
    bullets = [
        "LongCat-Flash-Chat 横向扩展：将 560B 模型扩展至万亿参数级，基于 MindSpeed-LLM，最少 2 个超节点（256 Die）即可完成训练。",
        "精度保持：扩展后模型在多个基准上复现原模型结果。",
        "开源贡献：",
        "   • PCL-Reasoner-V1/V1.5：模型权重、训练数据、代码已开源",
        "   • LLMEval：推理评估框架已开源至 Gitee 与启智社区",
        "未来方向：超长序列与超大规模 MoE 训练效率、Offline RL 与后训练流程、昇腾生态与开源社区深度融合。",
    ]
    _add_bullets(
        slide,
        bullets,
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.4),
        numbered=False,
        font_size=20,
    )
    _add_footer(slide, 14, TOTAL_PAGES)

    output_path = "/Users/jianzhengnie/work_dir/oh-my-claude-code/docs/昇腾生态与大模型训练实践.pptx"
    prs.save(output_path)
    print(f"PPT generated: {output_path}")


if __name__ == "__main__":
    main()
